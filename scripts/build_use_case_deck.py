#!/usr/bin/env python3
"""
Generate the Personal AI Framework (Skippy) use-case architecture deck.

Audience: a work colleague evaluating an edge NPU. They already understand
AI/ML; they don't care about UI polish. They want to know (1) what work the
AI/ML pipeline actually does, (2) how data flows through the system, (3) the
performance KPIs, (4) how it would perform on a target NPU.

Output: docs/personal-ai-use-cases.pptx (gitignored, private).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "personal-ai-use-cases.pptx"

# ── Theme ──
NAVY = RGBColor(0x0F, 0x19, 0x2E)
INK = RGBColor(0x11, 0x18, 0x27)
SURFACE = RGBColor(0x1A, 0x22, 0x3B)
ACCENT = RGBColor(0x63, 0x66, 0xF1)      # indigo
ACCENT2 = RGBColor(0x22, 0xC5, 0x5E)     # green (measured)
ACCENT3 = RGBColor(0xF5, 0x9E, 0x0B)     # amber (projected)
ACCENT4 = RGBColor(0xEF, 0x44, 0x44)     # red (gap)
TEXT = RGBColor(0xEA, 0xED, 0xF4)
MUTED = RGBColor(0x93, 0xA1, 0xB5)
FAINT = RGBColor(0x2A, 0x33, 0x4F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ──

def add_blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, mono=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if mono:
            run.font.name = "Consolas"
    return tb


def add_title(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.6),
             title, size=26, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                 subtitle, size=13, color=MUTED)


def add_box(slide, x, y, w, h, text, *, fill=SURFACE, border=ACCENT,
            size=12, bold=False, color=TEXT):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(80000); tf.margin_right = Emu(80000)
    tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold and i == 0
        run.font.color.rgb = color
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    ln = slide.shapes.add_connector(2, x1, y1, x2, y2)  # STRAIGHT
    ln.line.color.rgb = color
    ln.line.width = Pt(2)
    # Arrowhead
    fmt = ln.line
    fmt.fill.solid(); fmt.fill.fore_color.rgb = color
    xml = ln._element
    # Add arrow end element via raw XML
    from pptx.oxml.ns import qn
    from lxml import etree
    lnEl = xml.xpath(".//a:ln", namespaces={"a":"http://schemas.openxmlformats.org/drawingml/2006/main"})[0]
    tail = etree.SubElement(lnEl, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med"); tail.set("len", "med")


def add_bullets(slide, x, y, w, h, items, size=14, indent=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT
        p.space_after = Pt(4)
    return tb


def add_table(slide, x, y, w, h, headers, rows, *, header_fill=ACCENT,
              header_color=RGBColor(0xFF,0xFF,0xFF), band_fill=FAINT,
              font_size=11, highlight_rows=None):
    highlight_rows = highlight_rows or {}
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    # Header
    for c, h_text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Emu(40000); cell.margin_right = Emu(40000)
        cell.margin_top = Emu(20000); cell.margin_bottom = Emu(20000)
        tf = cell.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h_text
        r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = header_color
    # Rows
    for ri, row in enumerate(rows, start=1):
        fill = header_fill if ri in highlight_rows else (band_fill if ri % 2 == 0 else SURFACE)
        tcolor = RGBColor(0xFF,0xFF,0xFF) if ri in highlight_rows else TEXT
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c)
            cell.text = ""
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
            cell.margin_left = Emu(40000); cell.margin_right = Emu(40000)
            cell.margin_top = Emu(20000); cell.margin_bottom = Emu(20000)
            tf = cell.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(font_size); r.font.color.rgb = tcolor
            if ri in highlight_rows:
                r.font.bold = True
    return tbl


def add_footer(slide, n, total):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12),  Inches(0.3),
             f"Personal AI Framework (Skippy)  •  v5.9 architecture  •  {n}/{total}",
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ── Content ──

SLIDES = []


def slide1_title():
    s = add_blank()
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.2),
             "Personal AI Framework", size=48, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.7),
             "A fully-local AI assistant — use-case architecture",
             size=22, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             "Qwen 2.5 14B Instruct  |  61K-doc RAG  |  hybrid retrieval  |  tool-using agent  |  LoRA-trainable",
             size=14, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
             "Measured on RTX 5090 @ 32GB  •  Projected to 200-TOPS edge NPU",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
SLIDES.append(slide1_title)


def slide2_exec_summary():
    s = add_blank()
    add_title(s, "Executive summary",
              "What the framework does and why it matters for edge")
    add_bullets(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), [
        "Private AI assistant — runs 100% locally, no cloud API calls.",
        "Answers questions using your own corpus (61K+ documents: emails, transcripts, PDFs, source code).",
        "Retrieves context via hybrid search (semantic + BM25 + reranking), injects into Qwen 2.5 14B Instruct.",
        "Agentic: can call tools (read files, web search, send Gmail, create calendar events, schedule reminders, run scripts).",
        "Learns from use — 👍/👎 feedback excludes bad turns; LoRA retrains on your writing style.",
        "Multi-user with per-user data isolation (bcrypt auth, per-user SQLite + ChromaDB collections).",
        "Observable — Prometheus /metrics + Grafana dashboard for TTFT, throughput, RAG hits, tool calls.",
        "Why this deck matters: data flow + KPIs let us size a 200-TOPS edge NPU replacement for the RTX 5090 today.",
    ], size=14)
SLIDES.append(slide2_exec_summary)


def slide3_block_diagram():
    s = add_blank()
    add_title(s, "System block diagram",
              "Major components and the data paths between them")

    # Row 1: Input sources
    L = Inches(0.5); TOP = Inches(1.3); BW = Inches(2.2); BH = Inches(0.8); GAP = Inches(0.25)
    add_box(s, L, TOP, BW, BH, "Web UI\n(chat + voice)", fill=SURFACE, border=ACCENT)
    add_box(s, L + BW + GAP, TOP, BW, BH, "Upload\n(docs / audio / OCR)", fill=SURFACE, border=ACCENT)
    add_box(s, L + (BW+GAP)*2, TOP, BW, BH, "Mic (Whisper STT)\nin-browser", fill=SURFACE, border=ACCENT)
    add_box(s, L + (BW+GAP)*3, TOP, BW, BH, "Mobile / LAN\n(phone via IP)", fill=SURFACE, border=ACCENT)

    # Row 2: LLM Server (FastAPI) - spans
    LT = Inches(2.05)
    srv = add_box(s, L, LT, Inches(12.3), Inches(0.65),
                  "FastAPI LLM Server (auth middleware · per-user context · /metrics)",
                  fill=ACCENT, border=ACCENT, color=RGBColor(0xFF,0xFF,0xFF), bold=True, size=13)

    # Row 3: AI/ML subsystems
    T3 = Inches(2.85); H3 = Inches(1.7)
    add_box(s, L, T3, BW, H3,
            "Tool-detection\npass\n\n(short prompt,\nJSON out)",
            fill=SURFACE, border=ACCENT2, size=11)
    add_box(s, L + (BW+GAP), T3, BW, H3,
            "RAG retrieval\n\nChromaDB +\nBM25 + rerank\n(hybrid)",
            fill=SURFACE, border=ACCENT2, size=11)
    add_box(s, L + (BW+GAP)*2, T3, BW, H3,
            "Main LLM\ninference\n\nQwen 2.5 14B\nQ4_K_M GGUF\n(llama.cpp, CUDA)",
            fill=ACCENT2, border=ACCENT2, color=RGBColor(0x0F,0x19,0x2E), size=11, bold=True)
    add_box(s, L + (BW+GAP)*3, T3, BW, H3,
            "Memory\n(conversation)\n\nper-user\nChromaDB\ncollection",
            fill=SURFACE, border=ACCENT2, size=11)
    add_box(s, L + (BW+GAP)*4, T3, BW, H3,
            "Learned facts\n\nper-user\nChromaDB\ncollection",
            fill=SURFACE, border=ACCENT2, size=11)

    # Row 4: Storage + training
    T4 = Inches(4.75); H4 = Inches(1.3)
    add_box(s, L, T4, Inches(3.0), H4,
            "SQLite (per user)\nconversations · reminders · feedback (👍/👎) · settings",
            fill=SURFACE, border=ACCENT3, size=11)
    add_box(s, L + Inches(3.2), T4, Inches(3.0), H4,
            "Tool executor\nwrite_file · send_email · create_event · run_script · web_search · schedule_reminder",
            fill=SURFACE, border=ACCENT3, size=10)
    add_box(s, L + Inches(6.4), T4, Inches(3.0), H4,
            "LoRA training\nPEFT (rank 64) → GGUF merge → deploy\nSelective: skip 👎 and excluded convos",
            fill=SURFACE, border=ACCENT3, size=11)
    add_box(s, L + Inches(9.6), T4, Inches(3.2), H4,
            "Prometheus /metrics\nTTFT · tok/s · tool calls · 👍/👎 · RAG docs",
            fill=SURFACE, border=ACCENT3, size=11)

    # Legend
    add_text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.3),
             "Green = AI/ML inference  |  Amber = storage/data  |  Indigo = transport",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
SLIDES.append(slide3_block_diagram)


def slide4_inference_path():
    s = add_blank()
    add_title(s, "Data flow — single-turn inference (RAG)",
              "What happens between 'Enter' and the first token")

    steps = [
        ("User prompt", SURFACE, ACCENT),
        ("Query rewrite\n(LLM, 64 tok)", SURFACE, ACCENT2),
        ("Hybrid retrieval\nChromaDB + BM25 + rerank", SURFACE, ACCENT2),
        ("Top-k chunks\n(semantic + BM25 merged,\nreranked)", SURFACE, ACCENT2),
        ("Assemble prompt\nsystem + facts + memory\n+ RAG + history + user", SURFACE, ACCENT3),
        ("Main LLM inference\nQwen 2.5 14B\n(streaming)", ACCENT2, ACCENT2),
        ("Stream tokens\n→ Web UI (SSE)\n→ metrics recorded", SURFACE, ACCENT),
    ]

    L = Inches(0.5); T = Inches(1.6); BW = Inches(1.7); BH = Inches(1.4); GAP = Inches(0.1)
    for i, (txt, fill, border) in enumerate(steps):
        x = L + (BW + GAP) * i
        is_llm = fill == ACCENT2
        add_box(s, x, T, BW, BH, txt, fill=fill, border=border,
                size=10, bold=is_llm,
                color=RGBColor(0x0F,0x19,0x2E) if is_llm else TEXT)

    # KPIs below
    T2 = Inches(3.4)
    kpi_box = add_box(s, L, T2, Inches(12.3), Inches(1.0),
                      "Measured on RTX 5090  ·  Qwen 2.5 14B Q4_K_M  ·  16K ctx",
                      fill=INK, border=ACCENT, size=12, bold=True)
    add_text(s, L, Inches(4.5), Inches(12.3), Inches(2.6), [
        "• Query rewrite: 20–40 ms (small LLM call with 64 max tokens)",
        "• Hybrid retrieval: 60–120 ms (BM25 index + ChromaDB HNSW + cross-encoder reranker on top-20)",
        "• Prompt assembly: < 5 ms (pure Python string building)",
        "• TTFT (time to first token): 40–170 ms  — dominated by KV-cache warmup + prompt prefill",
        "• Throughput: 85–140 tok/s sustained  — bandwidth-bound at the model's 9 GB / context size",
        "• End-to-end for a 200-token answer: typically 1.5–3.0 s wall clock",
    ], size=13)
SLIDES.append(slide4_inference_path)


def slide5_agent_path():
    s = add_blank()
    add_title(s, "Data flow — agent tool pipeline",
              "Two-pass LLM: detect → (safe tools loop) → approve → execute → final generation")

    # Column 1: detection pass
    add_box(s, Inches(0.5), Inches(1.5), Inches(3.6), Inches(1.3),
            "Pass 1: Tool detection\nShort system prompt lists tools + few-shot examples.\nLLM emits <tool>...</tool> or nothing.\nmax_tokens=2048",
            fill=SURFACE, border=ACCENT2, size=11)
    # Column 2: loop
    add_box(s, Inches(0.5), Inches(3.0), Inches(3.6), Inches(2.3),
            "Pass 2a: Safe-tool loop (≤ 5 iterations)\n\nread_file · list_files · web_search · git_status · list_calendar_events\n\nEach result is fed back into detection prompt as <tool_output>…</tool_output>.\nLoop stops when the LLM emits no tool call.",
            fill=SURFACE, border=ACCENT2, size=11)
    # Column 3: confirm
    add_box(s, Inches(4.3), Inches(1.5), Inches(3.6), Inches(2.0),
            "Confirm tools\nwrite_file · run_script · send_email · create_calendar_event · schedule_reminder\n\nBreak the loop → surface pending-action card → wait for user approve/deny.",
            fill=SURFACE, border=ACCENT3, size=11)
    # Column 4: final gen
    add_box(s, Inches(4.3), Inches(3.7), Inches(3.6), Inches(1.6),
            "Pass 2b: Final generation\nAll tool results wrapped in <toolname_output>…</toolname_output> tags + scrubber to strip leaked scaffolding.",
            fill=ACCENT2, border=ACCENT2, color=RGBColor(0x0F,0x19,0x2E), size=11, bold=True)
    # Column 5: UI
    add_box(s, Inches(8.1), Inches(1.5), Inches(4.7), Inches(3.8),
            "UI effects\n• Collapsible 🔧 trace card per step\n• Pending-action card (param preview + diff)\n• Overwrite warning for existing files\n• 'Save as copy' sibling path option\n• Final streaming answer below traces",
            fill=SURFACE, border=ACCENT, size=11)

    add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.7), [
        "AI/ML cost breakdown (per user-visible turn, worst case):",
        "• 1 × detection pass (fast, short output)",
        "• ≤ 5 × safe-tool detections + executions (each ≈ 1 detection pass + 1 sandboxed exec)",
        "• 1 × RAG retrieval (reused from inference path)",
        "• 1 × final generation (full-length answer)",
        "Observed: adds 200–800 ms overhead on tool-using turns; tool execution itself is usually < 50 ms.",
    ], size=12)
SLIDES.append(slide5_agent_path)


def slide6_memory_rlhf():
    s = add_blank()
    add_title(s, "Data flow — memory + RLHF feedback loop",
              "How conversations flow into long-term memory and back into the model")

    # Flow: chat → SQLite → (auto) ingest → ChromaDB → retrieved as memory
    L = Inches(0.5); T = Inches(1.5); BW = Inches(2.0); BH = Inches(1.2); GAP = Inches(0.3)
    add_box(s, L, T, BW, BH, "Chat turn\n(user + assistant)", fill=SURFACE, border=ACCENT)
    add_box(s, L + (BW+GAP), T, BW, BH, "SQLite\nconversations.db\n(per-user)", fill=SURFACE, border=ACCENT3)
    add_box(s, L + (BW+GAP)*2, T, BW, BH, "Auto-ingest\n(on toggle)", fill=SURFACE, border=ACCENT2)
    add_box(s, L + (BW+GAP)*3, T, BW, BH, "ChromaDB\nmemory_<user>\ncollection", fill=SURFACE, border=ACCENT2)
    add_box(s, L + (BW+GAP)*4, T, BW, BH, "Next retrieval\nmemory hits\ninjected into ctx", fill=ACCENT2, border=ACCENT2, color=RGBColor(0x0F,0x19,0x2E), bold=True)

    # RLHF leg
    T2 = Inches(3.2)
    add_box(s, L, T2, BW, BH, "User rates\n👍 / 👎", fill=SURFACE, border=ACCENT)
    add_box(s, L + (BW+GAP), T2, BW, BH, "feedback table\n(in conversations.db)", fill=SURFACE, border=ACCENT3)
    add_box(s, L + (BW+GAP)*2, T2, BW, BH, "collect_training_data.py\nskips 👎 pairs\nskips excluded convos", fill=SURFACE, border=ACCENT2)
    add_box(s, L + (BW+GAP)*3, T2, BW, BH, "LoRA training\nPEFT rank 64\nbf16 · 3 epochs", fill=SURFACE, border=ACCENT2)
    add_box(s, L + (BW+GAP)*4, T2, BW, BH, "GGUF merge\n→ active model\nauto-restore", fill=ACCENT2, border=ACCENT2, color=RGBColor(0x0F,0x19,0x2E), bold=True)

    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), [
        "Why this matters for edge sizing:",
        "• Inference and memory retrieval are hot paths — every turn. These must run at interactive speed (< 3 s end-to-end).",
        "• Training is a background path — runs when the user triggers it, or on schedule. Can tolerate 2–3 h on device or offload to host.",
        "• RAG embeddings (ChromaDB) are cold storage — CPU-bound, not NPU-bound. 61K-doc index uses ~600 MB RAM + ~250 MB disk.",
        "• LoRA retraining is the NPU's stress-test case: Qwen 2.5 14B bf16 + rank-64 adapter ≈ 28 GB peak; usually offloaded on edge.",
    ], size=13)
SLIDES.append(slide6_memory_rlhf)


def slide7_kpis():
    s = add_blank()
    add_title(s, "Measured KPIs — what the RTX 5090 actually does",
              "Per-turn observability from Prometheus /metrics (v5.8+)")

    rows = [
        ("Time to first token (TTFT)",        "40–170 ms",       "p50 / p95",    "dominated by prompt prefill + KV warmup"),
        ("Throughput (streaming)",             "85–140 tok/s",     "sustained",    "bandwidth-bound at Q4_K_M · 9 GB model"),
        ("End-to-end, 200-token answer",       "1.5–3.0 s",         "wall",         "includes RAG + tool pass + generation"),
        ("RAG retrieval",                      "60–120 ms",         "k=5",          "hybrid (BM25 + HNSW + rerank top-20)"),
        ("Tool detection pass",                 "80–150 ms",         "64–2048 tok",  "separate short LLM call"),
        ("Tool execution (safe tools)",        "5–50 ms",           "local only",    "web_search: 300–800 ms (network)"),
        ("Per-user memory retrieval",          "40–80 ms",          "top-2",        "ChromaDB collection per user"),
        ("Peak VRAM (14B Q4_K_M, 16K ctx)",     "≈ 9.2 GB",           "steady",       "KV cache grows ~0.5 GB per 1K tokens"),
        ("Host CPU during inference",          "< 10%",              "single core",  "llama.cpp offloads all layers to GPU"),
        ("Docs in knowledge base",              "61,500+",            "ChromaDB",     "~600 MB RAM, 250 MB disk"),
        ("Model load time (cold)",             "20–30 s",            "GGUF mmap",    "one-time; hot restart re-uses page cache"),
    ]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0),
              ["Metric", "Measured", "How", "Notes"],
              rows, font_size=11)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             "All numbers from skippy_* Prometheus histograms on production usage.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
SLIDES.append(slide7_kpis)


def slide8_bw_math():
    s = add_blank()
    add_title(s, "Memory bandwidth — the first-order edge constraint",
              "LLM inference is memory-bandwidth-bound, not compute-bound")

    add_box(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(1.4),
            "For dense autoregressive models,\nevery token decode reads the full weight set at least once.\n\nBW_needed (GB/s) ≈ Model_size (GB) × tok/s_desired",
            fill=SURFACE, border=ACCENT2, size=14, bold=True)

    add_box(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(1.4),
            "Consequences for edge:\n• Big model = low tok/s unless BW is huge\n• Quantization (Q4 / INT4) halves the BW need\n• MoE only loads experts — lowers effective model size, but needs smart caching",
            fill=SURFACE, border=ACCENT3, size=12)

    # BW of various systems
    rows = [
        ("Vendor RTX 5090 (384-bit GDDR7)",       "1792 GB/s",    "reference",      "Measured"),
        ("Apple M3 Max (512-bit LPDDR5)",          "400 GB/s",     "reference",      "Measured"),
        ("128-bit LPDDR5x @ 8400 MT/s — peak",     "134.4 GB/s",   "peak",            "Theoretical"),
        ("128-bit LPDDR5x @ 8400 MT/s — 75% util", "100.8 GB/s",   "usable",          "Target NPU"),
        ("96-bit LPDDR5x @ 8400 MT/s — 75% util",   "75.6 GB/s",    "usable",          "Lower-bin NPU"),
    ]
    add_table(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(2.2),
              ["System", "BW", "Regime", "Source"],
              rows, font_size=11, highlight_rows={4})

    add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.9), [
        "Rule-of-thumb tok/s ceiling on the target (100.8 GB/s usable):",
        "• Qwen 2.5 3B Q4 (1.9 GB)  →  ~53 tok/s",
        "• Llama 2 7B Q4 (3.8 GB)  →  ~27 tok/s",
        "• Qwen 2.5 14B Q4 (8.7 GB)  →  ~12 tok/s   ← our current production model on 5090",
        "• Mixtral 8x7B Q4 (26 GB)  →  won't fit in RAM on most edge SKUs",
    ], size=13)
SLIDES.append(slide8_bw_math)


def slide9_model_comparison():
    s = add_blank()
    add_title(s, "Model catalog — dense vs MoE vs quantization",
              "Which model sizes make sense on an edge NPU")

    rows = [
        ("Qwen 2.5 1.5B Instruct",     "Dense", "0.9 GB",   "112 tok/s",  "Fits",       "Fine for narrow tasks; weak at reasoning"),
        ("Qwen 2.5 3B Instruct",        "Dense", "1.9 GB",   "53 tok/s",   "Fits",        "Good trade for edge QA"),
        ("Llama 2 7B Chat Q4",          "Dense", "3.8 GB",   "27 tok/s",   "Fits",        "Vendor NPU reference model"),
        ("Qwen 2.5 7B Instruct Q4",     "Dense", "4.4 GB",   "23 tok/s",   "Fits",        "Strong general model"),
        ("Qwen 2.5 14B Instruct Q4",    "Dense", "8.7 GB",   "12 tok/s",   "Fits (tight)", "Our production model; borderline interactive"),
        ("Qwen 2.5 32B Instruct Q4",    "Dense", "19 GB",    "5 tok/s",    "16GB: No",    "Unusable @ 100 GB/s BW"),
        ("Mixtral 8x7B Q4 (MoE 2/8)",    "MoE",   "26 GB",    "~15 tok/s", "16GB: No",    "Active params ≈ 13B; 26 GB RAM footprint"),
        ("DeepSeek-V2-Lite (MoE 2/64)", "MoE",   "16 GB",    "~84 tok/s",  "16GB: Tight", "Active params only 2.4B — edge-friendly MoE"),
    ]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.3),
              ["Model", "Type", "Q4 size", "tok/s ceiling", "16GB RAM?", "Notes"],
              rows, font_size=11, highlight_rows={5})
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             "tok/s ceilings assume 100.8 GB/s usable (128-bit LPDDR5x @ 75% util). Actual numbers depend on kernel efficiency.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
SLIDES.append(slide9_model_comparison)


def slide10_target_npu():
    s = add_blank()
    add_title(s, "Target NPU case study",
              "200 TOPS · 128-bit LPDDR5x @ 8.4 GT/s · 75% utilization")

    add_box(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.6),
            "Specs as given",
            fill=SURFACE, border=ACCENT, size=12, bold=True)
    add_text(s, Inches(0.7), Inches(1.9), Inches(5.4), Inches(2.0), [
        "• INT8 compute: 200 TOPS",
        "• Memory bus: 128-bit LPDDR5x",
        "• Data rate: 8400 MT/s (= 8.4 GT/s)",
        "• Utilization budget: 75%",
        "• Peak BW = 128/8 × 8.4 = 134.4 GB/s",
        "• Usable BW = 134.4 × 0.75 = 100.8 GB/s",
    ], size=13)

    add_box(s, Inches(6.6), Inches(1.4), Inches(6.3), Inches(2.6),
            "Compute vs bandwidth",
            fill=SURFACE, border=ACCENT2, size=12, bold=True)
    add_text(s, Inches(6.8), Inches(1.9), Inches(5.9), Inches(2.0), [
        "• Llama 2 7B Q4: ~3.8 GB weights",
        "• Decode reads ~3.8 GB/token → 100.8/3.8 ≈ 27 tok/s",
        "• Compute for 7B at Q4/INT4: ~14 GOPs/token (attention + FFN)",
        "• At 200 TOPS INT8 ceiling → compute ceiling: 14000 tok/s",
        "• Conclusion: bandwidth-bound by ~520×, not compute-bound.",
    ], size=13)

    # Bottom chart: different utilization scenarios
    rows = [
        ("Conservative (50% util)",   "67.2 GB/s",   "17.7 tok/s",    "Safe design margin"),
        ("Moderate (60% util)",        "80.64 GB/s",  "21.2 tok/s",    "Conservative vendor claim"),
        ("Target (75% util)",          "100.8 GB/s",   "26.5 tok/s",    "This deck's working number"),
        ("Aggressive (85% util)",      "114.2 GB/s",   "30.0 tok/s",    "Requires excellent memory controller"),
    ]
    add_table(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.4),
              ["Utilization scenario", "Usable BW", "Llama 2 7B Q4 tok/s", "Comment"],
              rows, font_size=11, highlight_rows={3})
SLIDES.append(slide10_target_npu)


def slide_moe_memory_model():
    s = add_blank()
    add_title(s, "MoE vs dense — the edge trade-off",
              "Why a sparse mixture-of-experts model changes the bandwidth math")

    add_box(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(2.4),
            "Dense model (Qwen 2.5 14B)",
            fill=SURFACE, border=ACCENT2, size=13, bold=True)
    add_text(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(1.9), [
        "• Every token decode reads ALL weights once",
        "• BW_needed = full_model_size × tok/s",
        "• RAM footprint = full_model_size",
        "• Capacity = active_params (same as total)",
        "• Simple: one bandwidth number, one RAM number",
    ], size=13)

    add_box(s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(2.4),
            "MoE model (e.g. Mixtral 8x7B, 2-of-8)",
            fill=SURFACE, border=ACCENT3, size=13, bold=True)
    add_text(s, Inches(6.9), Inches(1.95), Inches(5.7), Inches(1.9), [
        "• Every token routes to k-of-N experts (e.g. 2-of-8)",
        "• BW_needed = active_expert_size × tok/s",
        "• RAM footprint = total_model_size (all experts resident)",
        "• Capacity = total_params (richer knowledge)",
        "• Edge wins: low BW per token, but high RAM ceiling",
    ], size=13)

    # Formula + example
    add_box(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.7),
            "The MoE edge appeal in one line",
            fill=INK, border=ACCENT, size=13, bold=True)
    add_text(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(2.0), [
        "BW_needed (GB/s) ≈ active_expert_size (GB) × tok/s_desired",
        "",
        "Mixtral 8x7B Q4: total 26 GB, active (2/8 × shared) ≈ 6.5 GB per token.",
        "At 100.8 GB/s usable → ~15 tok/s decode ceiling — higher than 14B dense (12 tok/s) at equivalent quality.",
        "But the 26 GB RAM requirement excludes most edge SKUs (16 GB).",
        "",
        "Qwen 3 30B-A3B (3B active): total 16 GB, active ~1.5 GB — measured 37.85 TPS on Edge NPU 2.",
        "DeepSeek-V2-Lite (2/64, 2.4B active): total 16 GB, active ~1.2 GB — ceiling ~84 tok/s.",
    ], size=13, mono=False)
SLIDES.append(slide_moe_memory_model)


def slide_moe_vs_dense_on_target():
    s = add_blank()
    add_title(s, "MoE vs dense on the target NPU",
              "Same 200 TOPS · 80.64 GB/s usable — which model family wins?")

    rows = [
        # (Model, Type, Total Q4 RAM, Active per-tok, BW-bound tok/s, 16GB edge?, Quality class)
        ("Qwen 2.5 3B",                 "Dense",  "1.9 GB",    "1.9 GB",    "53 tok/s",      "Fits",           "Small"),
        ("Llama 2 7B",                  "Dense",  "3.8 GB",    "3.8 GB",    "27 tok/s",      "Fits",           "Mid"),
        ("Qwen 2.5 7B",                 "Dense",  "4.4 GB",    "4.4 GB",    "23 tok/s",      "Fits",           "Mid"),
        ("Qwen 2.5 14B",                "Dense",  "8.7 GB",    "8.7 GB",    "12 tok/s",      "Fits (tight)",   "Large ← our current"),
        ("Qwen 3 30B-A3B",              "MoE",    "16 GB",     "1.5 GB",    "~67 tok/s ceiling  ·  37.85 TPS measured (Edge NPU 2)", "Tight",          "Large ★"),
        ("DeepSeek-V2-Lite (2/64)",      "MoE",    "16 GB",     "1.2 GB",    "~84 tok/s",     "Tight (swap?)",  "Mid-Large"),
        ("Mixtral 8x7B (2/8)",          "MoE",    "26 GB",     "6.5 GB",    "~15 tok/s",     "Does NOT fit",   "Large"),
        ("Qwen 1.5 MoE-A2.7B (4/60)",    "MoE",    "14.3 GB",   "1.3 GB",    "~77 tok/s",     "Tight",           "Mid"),
    ]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.2),
              ["Model", "Type", "Q4 total RAM", "Active / token", "tok/s ceiling", "16 GB edge?", "Quality tier"],
              rows, font_size=11, highlight_rows={4})  # highlight Qwen 3 30B-A3B (measured)

    add_text(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.6), [
        "Takeaways (★ = measured, not just theoretical):",
        "• Qwen 3 30B-A3B on Edge NPU 2: TTFT 0.351 s @ 1K prompt, 37.85 TPS — 56% of BW ceiling, realistic vendor number.",
        "• On a 16 GB edge SKU, Mixtral 8x7B is out (26 GB total). Qwen 3 30B-A3B + DeepSeek-V2-Lite both squeeze in.",
        "• MoE at 3B active delivers ~3× the tok/s of 14B dense on the same bandwidth, with larger effective capacity.",
        "• Training MoE on-device is not viable at 16 GB — offload to host (5090 can do Q4 QLoRA with gradient checkpointing, ~4–8 h).",
    ], size=12)
SLIDES.append(slide_moe_vs_dense_on_target)


def slide11_vendor_claim():
    s = add_blank()
    add_title(s, "Vendor claim reconciliation",
              "60 tok/s Llama 2 7B — what would have to be true")

    add_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.6),
            "The gap",
            fill=SURFACE, border=ACCENT4, size=12, bold=True)
    add_text(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(2.0), [
        "• Physics ceiling @ 75% util: 26.5 tok/s",
        "• Vendor claim: 60 tok/s",
        "• Gap: 2.26× over the bandwidth bound",
        "• That still rules out a plain FP16 / Q4 decode loop",
    ], size=13)

    add_box(s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(2.6),
            "Paths to 60 tok/s",
            fill=SURFACE, border=ACCENT2, size=12, bold=True)
    add_text(s, Inches(6.9), Inches(1.9), Inches(5.7), Inches(2.0), [
        "• INT4 weight-only → halves BW need → ~53 tok/s",
        "• Speculative decoding (draft model) → 1.3–1.5× wins",
        "• Either alone gets close; combined exceeds 60 tok/s",
        "• Flash-attention v2 kernels → marginal for decode",
        "• Or: benchmark ran at 90%+ memory utilization",
    ], size=13)

    add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(3.0), [
        "What to ask the vendor:",
        "• Quantization scheme (INT4 weight-only? INT8 activations? AWQ / GPTQ?)",
        "• Prompt length used in the benchmark (typical: 128; typical for RAG: 4–8 K)",
        "• Is the 60 tok/s single-stream decode, or batched throughput?",
        "• Speculative decoding on/off? If on, what draft model?",
        "• Sustained vs peak? Over how many generated tokens?",
        "• What's the tok/s for a 7B model with an 8 K-token RAG prompt? (Our real workload)",
    ], size=13)
SLIDES.append(slide11_vendor_claim)


def slide12_workload_fit():
    s = add_blank()
    add_title(s, "Which workload fits on the target",
              "Mapping Skippy's current use cases to the 200-TOPS / 80 GB/s NPU")

    rows = [
        ("Single-user chat, 3B dense, 4K ctx",      "Fits · 53 tok/s",    "Ideal edge case"),
        ("Single-user chat, 7B dense, 4K ctx",      "Fits · 27 tok/s",    "Comfortably interactive"),
        ("Single-user chat, 14B dense, 16K ctx",    "Tight · 12 tok/s",   "Our current model; acceptable on this NPU"),
        ("Single-user chat, Qwen3 30B-A3B, 4K",     "Fits · 37.85 TPS ★", "Measured on Edge NPU 2 — strong edge MoE"),
        ("Hybrid RAG (+0.1 s retrieval)",            "Negligible overhead",  "CPU-side; not on NPU"),
        ("Tool detection pass",                       "Adds 2nd small LLM call",  "Same model, 64–2048 tok — cheap"),
        ("Agent loop (≤ 5 safe tools)",                "Adds 5 × detection ≈ 0.5 s",  "User-visible"),
        ("Whisper STT (voice input)",                    "Runs on CPU/GPU, ~4× realtime",  "Not on NPU typically"),
        ("OCR (Tesseract) / PDF ingest",                 "CPU-bound, minutes",             "Batch/background"),
        ("LoRA retraining (rank 64, 3 epoch)",           "Needs ≈ 28 GB RAM + fp16 compute",  "OFFLOAD — not on this NPU"),
        ("Multi-user, 3 concurrent sessions",            "Needs batching or queue",            "llama.cpp is single-stream; contention"),
        ("Prometheus /metrics endpoint",                  "Negligible",                          "Plain CPU"),
    ]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
              ["Workload", "Verdict", "Note"],
              rows, font_size=11, highlight_rows={2})
SLIDES.append(slide12_workload_fit)


def slide13_platform_sizing():
    s = add_blank()
    add_title(s, "Platform sizing — reference points",
              "Where the target NPU lands on the spectrum")

    rows = [
        ("NVIDIA RTX 5090 (desktop)",        "209 TOPS (INT8)",   "1792 GB/s",   "Qwen 2.5 14B @ 85–140 tok/s",              "450 W"),
        ("NVIDIA Jetson AGX Orin 64GB",       "275 TOPS (INT8)",    "204 GB/s",    "Qwen 2.5 7B @ 25–40 tok/s",                 "15–60 W"),
        ("Apple M3 Max (128GB)",              "~18 TFLOPS GPU",      "400 GB/s",    "Qwen 2.5 14B @ 30–50 tok/s",                 "~70 W"),
        ("Target NPU (this deck, 75% util)",   "200 TOPS (INT8)",     "100.8 GB/s",  "Qwen 3 30B-A3B @ 37.85 TPS ★ (measured)",    "TBD"),
        ("Mobile SoC (reference)",             "~45 TOPS",             "~60 GB/s",     "Qwen 2.5 3B @ ~10 tok/s",                     "3–5 W"),
    ]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.0),
              ["Platform", "Compute", "BW", "Typical Skippy model + tok/s", "Power"],
              rows, font_size=11, highlight_rows={4})

    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), [
        "Observations:",
        "• The target NPU has GPU-class compute (200 TOPS) but phone-class memory bandwidth (100.8 GB/s usable @ 75% util).",
        "• For interactive 7B dense chat this is good (~27 tok/s); for 14B dense it's borderline interactive (~12 tok/s).",
        "• Qwen 3 30B-A3B MoE measured at 37.85 TPS on Edge NPU 2 — the strongest option at 16 GB RAM.",
        "• Jetson AGX Orin is closest in spirit but with 2× the BW — useful as a benchmark baseline.",
        "• Mobile SoCs with ~60 GB/s BW land at the 3B model tier for a smooth UX.",
    ], size=13)
SLIDES.append(slide13_platform_sizing)


def slide14_takeaways():
    s = add_blank()
    add_title(s, "Key takeaways",
              "What the colleague should walk away with")

    add_bullets(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6), [
        "Skippy is memory-bandwidth-bound, not compute-bound — 200 TOPS is comfortably over-provisioned; 100.8 GB/s usable BW (75% util) is the real constraint.",
        "For dense models at 100.8 GB/s: Llama 2 7B Q4 tops out at ~27 tok/s; Qwen 2.5 14B (our current production) sits at ~12 tok/s — borderline interactive on this NPU.",
        "MoE with small active-params wins the edge: Qwen 3 30B-A3B delivers measured 37.85 TPS on Edge NPU 2 with TTFT 0.351 s (1K prompt) — ~3× the throughput of 14B dense at larger effective capacity.",
        "Practical recommendation for this NPU class: Qwen 3 30B-A3B if the 16 GB RAM budget fits total model, Qwen 2.5 3B–7B dense as safe fallback.",
        "Mixtral 8x7B is OUT for 16 GB edge SKUs (26 GB total). Worth revisiting only if the target ships with 32 GB+.",
        "The agent/RAG/memory/RLHF machinery around the LLM is cheap compared to generation — don't let it distract from the bandwidth math.",
        "Training stays offload — LoRA retraining ≈ 28 GB peak VRAM; run on the host, not the edge NPU. Q4 QLoRA on MoE is possible on a 5090 with gradient checkpointing (~4–8 h).",
        "Observability matters — Prometheus /metrics lets us validate any NPU claim against the same histograms used in production.",
    ], size=13)
SLIDES.append(slide14_takeaways)


# ── Build ──

def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    for gen in SLIDES:
        gen()
    # Page-number footers now that we know the total.
    total = len(SLIDES)
    for i, slide in enumerate(prs.slides, start=1):
        add_footer(slide, i, total)
    prs.save(OUT)
    print(f"[build_use_case_deck] Wrote {total} slides to {OUT} ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
